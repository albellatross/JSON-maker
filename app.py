import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import io
import zipfile
import random
import urllib.parse
import re
import base64

# ================= 🎨 1. DESIGN TOKENS & UI CONFIGURATION =================
MY_DESIGN_TOKENS = {
    "bg_color": "#FFF8F3",        # 极浅的暖米色背景
    "surface_color": "#FFFFFF",   # 纯白卡片背景
    "text_primary": "#311F10",    # 深棕色主文字
    "text_secondary": "#6B5A50",  # 浅棕色次要文字
    "accent_color": "#4A3A32",    # 强调色（按钮、激活状态）
    "accent_hover": "#635147",    # 按钮悬停色
    "border_color": "#E8DED5",    # 柔和的边框色
    "font_family": "'Segoe UI', 'Helvetica Neue', sans-serif"
}

def inject_layout_css(tokens):
    css = f"""
    <style>
        /* === 1. 全局布局重置 (No-Scroll Core) === */
        .stApp {{
            background-color: {tokens['bg_color']};
            font-family: {tokens['font_family']};
            color: {tokens['text_primary']};
            overflow: hidden !important; /* 禁止整个网页滚动 */
        }}
        
        /* 隐藏顶部 Header 和 Footer */
        header, footer, [data-testid="stHeader"] {{ display: none !important; }}
        
        /* 极致压缩页面边距，最大化利用屏幕空间 */
        .block-container {{
            padding-top: 1rem !important;
            padding-bottom: 0 !important;
            padding-left: 1.5rem !important;
            padding-right: 1.5rem !important;
            max-width: 100% !important;
            height: 100vh !important;
        }}

        h1, h2, h3, h4 {{ margin: 0 !important; padding: 0 !important; color: {tokens['text_primary']} !important; }}

        /* === 2. 顶部菜单 (Tabs) 美化 === */
        /* Tab 容器 */
        .stTabs [data-baseweb="tab-list"] {{
            gap: 8px;
            background-color: transparent;
            padding-bottom: 0;
            margin-bottom: 1rem;
            border-bottom: 1px solid {tokens['border_color']};
        }}

        /* 单个 Tab 按钮 */
        .stTabs [data-baseweb="tab"] {{
            height: 40px;
            white-space: pre-wrap;
            border-radius: 8px 8px 0 0;
            color: {tokens['text_secondary']};
            font-weight: 600;
            border: none;
            background-color: transparent;
            padding: 0 20px;
            transition: all 0.2s ease;
        }}

        /* 选中状态的 Tab */
        .stTabs [aria-selected="true"] {{
            background-color: {tokens['surface_color']} !important;
            color: {tokens['accent_color']} !important;
            border: 1px solid {tokens['border_color']} !important;
            border-bottom: 1px solid {tokens['surface_color']} !important; /* 遮住底部分割线 */
            box-shadow: 0 -2px 5px rgba(0,0,0,0.02);
        }}
        
        /* 去掉 Tab 选中时的红线 */
        .stTabs [data-baseweb="tab-highlight"] {{ display: none; }}

        /* === 3. 左右分栏布局 (核心高度控制) === */
        
        /* 左侧面板：图片展示区 */
        .left-panel {{
            /* 动态计算高度：总高 - 顶部Tab高度 - 边距 */
            height: calc(100vh - 100px); 
            background-color: #F2EBE6; /* 稍微深一点的背景衬托图片 */
            border-radius: 16px;
            display: flex;
            justify-content: center;
            align_items: center;
            overflow: hidden;
            border: 1px solid {tokens['border_color']};
            position: relative;
        }}
        
        .left-panel img {{
            max-width: 95%;
            max-height: 95%;
            width: auto;
            height: auto;
            object-fit: contain; /* 关键：保持比例不裁切 */
            filter: drop-shadow(0 8px 16px rgba(0,0,0,0.08));
        }}

        /* 右侧面板：独立滚动区 */
        .right-scroll-area {{
            height: calc(100vh - 100px); /* 与左侧等高 */
            overflow-y: auto; /* 仅允许此处滚动 */
            padding-right: 10px;
            padding-left: 5px;
            padding-bottom: 40px; /* 给底部留点呼吸空间 */
        }}
        
        /* 隐藏/美化右侧滚动条 */
        .right-scroll-area::-webkit-scrollbar {{ width: 6px; }}
        .right-scroll-area::-webkit-scrollbar-track {{ background: transparent; }}
        .right-scroll-area::-webkit-scrollbar-thumb {{ background-color: #DCCFC6; border-radius: 3px; }}
        .right-scroll-area::-webkit-scrollbar-thumb:hover {{ background-color: #C4B4A8; }}

        /* === 4. 组件微调 === */
        
        /* 卡片 */
        [data-testid="stVerticalBlock"] > [style*="flex-direction: column;"] > [data-testid="stVerticalBlock"] {{
            background-color: {tokens['surface_color']};
            border-radius: 12px;
            padding: 1.2rem;
            border: 1px solid {tokens['border_color']};
            box-shadow: 0 2px 4px rgba(0,0,0,0.02);
            gap: 0.8rem;
        }}

        /* 输入框 */
        .stTextArea textarea, .stTextInput input {{
            font-size: 13px;
            border-radius: 8px;
            border: 1px solid {tokens['border_color']};
            background-color: #FAFAFA;
        }}
        .stTextArea textarea:focus, .stTextInput input:focus {{
            border-color: {tokens['accent_color']};
            box-shadow: 0 0 0 1px {tokens['accent_color']};
        }}

        /* 按钮 */
        .stButton button {{
            border-radius: 8px !important;
            font-weight: 600 !important;
            font-size: 14px !important;
            transition: all 0.2s;
        }}
        
        /* 主按钮 (Primary) */
        div[data-testid="stButton"] > button[kind="primary"] {{ 
            background-color: {tokens['accent_color']} !important; 
            color: #FFFFFF !important; 
            border: none !important;
            padding: 0.5rem 1rem;
        }}
        div[data-testid="stButton"] > button[kind="primary"]:hover {{ 
            background-color: {tokens['accent_hover']} !important;
            box-shadow: 0 4px 12px rgba(74, 58, 50, 0.2);
        }}

        /* 图片圆角 */
        img {{ border-radius: 8px !important; }}
        
        /* 隐藏链接图标 */
        .css-1v0mbdj a {{ display: none; }}
        
        /* 进度条 */
        .stProgress > div > div > div > div {{ background-color: {tokens['accent_color']}; }}
        
        /* 修正 Streamlit 默认的 margin */
        .element-container {{ margin-bottom: 0.5rem !important; }}
    </style>
    """
    st.markdown(css, unsafe_allow_html=True)

# ================= 2. DATA LISTS =================

REMIX_LIST_EN = [
    {"label": "Want a wider view?", "prompt": "Create an expanded image with extended space."},
    {"label": "Want to zoom in?", "prompt": "Create a micro-detail close-up variant of this image."},
    {"label": "Try paper cut style?", "prompt": "Remake this image in a modern paper cut style with layered colors and soft shadows."},
    {"label": "Make this embroidery style?", "prompt": "Remake this image in a textile embroidery style with visible stitched threads."},
    {"label": "Change to Pixel Art", "prompt": "Create this picture as a retro pixel art, with nostalgic detail and game shading."},
    {"label": "Apply Glitch Effect", "prompt": "Remake this image as a glitch digital art, with pixel splits and cyberpunk noise."},
    {"label": "Change to Watercolor", "prompt": "Create this picture as a watercolor painting."},
    {"label": "Change to Impressionism", "prompt": "Create this picture as an Impressionist painting, with loose brushwork, luminous color, and fleeting light."},
    {"label": "Draw with colored pencil?", "prompt": "Remake this image as a colored pencil drawing."},
    {"label": "Try fine-line style?", "prompt": "Remake this image as a Chinese Gongbi painting with precise outlines, soft washes, and detailed forms."},
    {"label": "Try Chinese paper cut style?", "prompt": "Remake this image as a Chinese paper cut, with red silhouettes, cultural motifs, and symmetrical patterns."},
    {"label": "Try Ukiyo-e style?", "prompt": "Remake this image as a Japanese Ukiyo-e, with woodblock texture, flat colors, and flowing lines."},
    {"label": "Make this a portrait?", "prompt": "Remake this image as a photo portrait, with natural light, and shallow depth."},
    {"label": "Make this stained glass?", "prompt": "Remake this image as a stained glass design with colorful panes, bold outlines, and glowing light."},
    {"label": "Try silkscreen style?", "prompt": "Remake this image as a silkscreen print."},
    {"label": "Make this anime?", "prompt": "Remake this image as an anime illustration with expressive light and a dynamic layout."},
    {"label": "Add sepia tone?", "prompt": "Remake this image as a sepia-toned memory with aged paper texture."},
    {"label": "Make this pop art?", "prompt": "Remake this image as a high-saturation pop art, with bold blocks and hues."},
    {"label": "Make this a gradient mesh?", "prompt": "Remake this image as a gradient mesh, blending colors seamlessly across the composition."},
    {"label": "Make this a 3D figure?", "prompt": "Remake this image as a photorealistic 3D render of a collectible figure, made of real materials like resin or plastic with cinematic lighting, studio backdrop, and ultra-fine modeling detail."},
    {"label": "Try duotone colors?", "prompt": "Remake this image as a duotone image."},
    {"label": "Make this monochrome?", "prompt": "Remake this image as a monochrome image."},
    {"label": "Add neon lighting?", "prompt": "Remake this image as a neon-lit scene with vibrant color contrasts."},
    {"label": "Make this mechanical?", "prompt": "Create a mechanical version of the subject with exposed gears, metallic joints, and precise components."},
    {"label": "Make this crystal?", "prompt": "Remake this image to be in an iridescent fantasy realm with the subject as translucent glass or crystal, glowing and refracted."}
]

def get_random_remix(): return random.choice(REMIX_LIST_EN)

def randomize_callback(index, session_key_root, current_id_val):
    new_remix = get_random_remix()
    st.session_state[session_key_root][index] = new_remix
    st.session_state[f"l_{current_id_val}_{index}"] = new_remix['label']
    st.session_state[f"p_{current_id_val}_{index}"] = new_remix['prompt']

def parse_bulk_remix_text(raw_text):
    if not raw_text.strip(): return []
    ACTION_KEYWORDS = ("create", "change", "recreate", "replace", "generate", "make", "transform", "add", "switch", "use", "apply", "convert", "turn")
    lines = raw_text.split('\n')
    parsed_items = []
    processed_indices = set() 
    def clean_line_start(s): return re.sub(r'^[\d\.\-\*\s]+', '', s).strip()

    for i, line in enumerate(lines):
        clean_current = clean_line_start(line)
        lower_current = clean_current.lower()
        is_prompt_start = lower_current.startswith(ACTION_KEYWORDS)
        inline_split = line.split(":", 1)
        has_inline_title = len(inline_split) > 1 and clean_line_start(inline_split[1]).lower().startswith(ACTION_KEYWORDS)

        if is_prompt_start or has_inline_title:
            title = "Remix Option"
            prompt_text = ""
            if has_inline_title:
                title = clean_line_start(inline_split[0])
                prompt_text = clean_line_start(inline_split[1])
                processed_indices.add(i)
            else:
                prompt_text = clean_current
                processed_indices.add(i)
                k = i - 1
                while k >= 0:
                    prev_line = lines[k].strip()
                    if prev_line and k not in processed_indices:
                        title = clean_line_start(prev_line)
                        processed_indices.add(k)
                        break
                    k -= 1
            j = i + 1
            while j < len(lines):
                next_line = lines[j].strip()
                if not next_line: j+=1; continue
                clean_next = clean_line_start(next_line)
                if clean_next.lower().startswith(ACTION_KEYWORDS): break
                if j + 1 < len(lines):
                    clean_next_next = clean_line_start(lines[j+1])
                    if clean_next_next.lower().startswith(ACTION_KEYWORDS): break
                prompt_text += " " + next_line
                processed_indices.add(j)
                j += 1
            parsed_items.append({"label": title, "prompt": prompt_text})
    return parsed_items

def batch_parse_callback(session_key, current_id_val):
    batch_text = st.session_state.get("batch_input_area", "")
    parsed_items = parse_bulk_remix_text(batch_text)
    if parsed_items:
        final_items = parsed_items[:3]
        while len(final_items) < 3:
            final_items.append(get_random_remix())
        st.session_state[session_key] = final_items
        for idx, item in enumerate(final_items):
            st.session_state[f"l_{current_id_val}_{idx}"] = item['label']
            st.session_state[f"p_{current_id_val}_{idx}"] = item['prompt']
        st.session_state["batch_input_area"] = ""
        st.session_state["_parse_success"] = True
        st.session_state["_parsed_count"] = len(parsed_items)
    else:
        st.session_state["_parse_error"] = True

# --- File Operations ---

def process_ppt_file(uploaded_file, start_id):
    uploaded_file.seek(0)
    try:
        prs = Presentation(uploaded_file)
    except zipfile.BadZipFile:
        raise ValueError("File is not a valid .pptx file.")
    except Exception as e:
        raise ValueError(f"Error reading PPT: {str(e)}")

    current_id = int(start_id)
    extracted_data = []
    image_storage = {}
    for index, slide in enumerate(prs.slides):
        slide_info = {"id": str(current_id), "original_prompt_text": "", "image_filename": ""}
        found_image = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not found_image:
                img_name = f"{current_id}.png"
                image_storage[img_name] = shape.image.blob
                slide_info["image_filename"] = img_name
                found_image = True
            if shape.has_text_frame:
                text = shape.text.strip()
                if len(text) > 10 and len(text) > len(slide_info["original_prompt_text"]):
                    slide_info["original_prompt_text"] = text
        if found_image:
            extracted_data.append(slide_info)
            current_id += 1
    return extracted_data, image_storage

def create_final_zip(processed_jsons, image_storage):
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        combined_list = []
        sorted_keys = sorted(processed_jsons.keys(), key=lambda x: int(processed_jsons[x]['id']))
        for key in sorted_keys:
            item = processed_jsons[key]
            clean_item = {
                "id": item.get("id"),
                "prompt": item.get("prompt"),
                "remixSuggestions": [
                    {"label": r.get("label"), "prompt": r.get("prompt")} for r in item.get("remixSuggestions", [])
                ]
            }
            combined_list.append(clean_item)
            target_img_name = f"{item.get('id')}.png"
            if target_img_name in image_storage:
                zip_file.writestr(f"images/{target_img_name}", image_storage[target_img_name])
        json_str = json.dumps(combined_list, indent=4, ensure_ascii=False)
        zip_file.writestr("dataset.json", json_str)
    return zip_buffer

def renumber_json_ids(json_file, start_num):
    try:
        content = json.load(json_file)
        if not isinstance(content, list): return None, "Error: JSON root must be a list []"
        counter = int(start_num)
        for item in content:
            if 'id' in item:
                item['id'] = str(counter)
                counter += 1
        return json.dumps(content, indent=4, ensure_ascii=False), None
    except Exception as e:
        return None, f"Error: {str(e)}"

def extract_images_from_ppt(uploaded_file, start_id):
    uploaded_file.seek(0)
    try:
        prs = Presentation(uploaded_file)
    except Exception as e:
        return None, f"Error: {str(e)}"
    zip_buffer = io.BytesIO()
    current_id = int(start_id)
    count = 0
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_name = f"{current_id}.png"
                    img_bytes = shape.image.blob
                    zip_file.writestr(img_name, img_bytes)
                    current_id += 1
                    count += 1
    return zip_buffer, count

# ================= 3. MAIN UI =================
st.set_page_config(page_title="Remix Studio", layout="wide", page_icon="🧶")
inject_layout_css(MY_DESIGN_TOKENS)

# Tabs
tab_main, tab_fix, tab_extract = st.tabs(["🧶 Remix Editor", "🔢 JSON ID Fixer", "🖼️ PPT Image Extractor"])

# ================= TAB 1: REMIX EDITOR =================
with tab_main:
    if 'data' not in st.session_state: st.session_state.data = []
    if 'images' not in st.session_state: st.session_state.images = {}
    if 'processed_results' not in st.session_state: st.session_state.processed_results = {}
    if 'current_idx' not in st.session_state: st.session_state.current_idx = 0

    if not st.session_state.data:
        st.markdown("<br><br>", unsafe_allow_html=True)
        c1, c2, c3 = st.columns([1, 2, 1])
        with c2:
            st.markdown(f"<div style='text-align: center; margin-bottom:20px;'><h1>🧶 Remix Studio</h1><p style='color:#64748B'>Upload PPTX • Edit • Export</p></div>", unsafe_allow_html=True)
            with st.container(border=True):
                uploaded_ppt = st.file_uploader("Upload Presentation", type=["pptx"])
                start_id = st.number_input("Start ID", value=453, step=1)
                st.markdown("<br>", unsafe_allow_html=True)
                if uploaded_ppt:
                    if st.button("🚀 Load Slides", type="primary", use_container_width=True):
                        with st.spinner("Processing..."):
                            try:
                                data, images = process_ppt_file(uploaded_ppt, start_id)
                                if not data: st.error("No valid slides found.")
                                else:
                                    st.session_state.data = data
                                    st.session_state.images = images
                                    st.session_state.current_idx = 0
                                    st.rerun()
                            except ValueError as ve: st.error(str(ve))
                            except Exception as e: st.error(f"Error: {e}")

    else:
        item = st.session_state.data[st.session_state.current_idx]
        current_id = item['id']
        img_name = item['image_filename']

        # 布局：左侧 45%，右侧 55%
        col_left, col_right = st.columns([1.2, 1.5], gap="medium")

        # === LEFT (Image Panel) ===
        with col_left:
            st.markdown(f"#### ID {current_id}")
            if img_name in st.session_state.images:
                b64_img = base64.b64encode(st.session_state.images[img_name]).decode()
                st.markdown(f"""
                <div class="left-panel">
                    <img src="data:image/png;base64,{b64_img}" />
                </div>
                """, unsafe_allow_html=True)
            else:
                st.error("Image missing")

        # === RIGHT (Editor Panel) ===
        with col_right:
            # 1. Top Bar: Progress & Export
            c_top1, c_top2 = st.columns([3, 1])
            with c_top1:
                total = len(st.session_state.data)
                done = len(st.session_state.processed_results)
                st.progress(done / total if total > 0 else 0)
                st.caption(f"Progress: {done} / {total} (ID: {current_id})")
            
            with c_top2:
                with st.popover("⚙️ Export"):
                    if done >= 0:
                        export_data = st.session_state.processed_results.copy()
                        curr_main = st.session_state.get(f"m_{current_id}", "")
                        curr_remixes = []
                        session_key = f"remix_{current_id}"
                        if session_key in st.session_state:
                            raw_remixes = st.session_state[session_key]
                            for idx in range(len(raw_remixes)):
                                l_val = st.session_state.get(f"l_{current_id}_{idx}", raw_remixes[idx]['label'])
                                p_val = st.session_state.get(f"p_{current_id}_{idx}", raw_remixes[idx]['prompt'])
                                curr_remixes.append({"label": l_val, "prompt": p_val})
                        
                        export_data[f"{current_id}.json"] = {
                            "id": current_id,
                            "prompt": curr_main,
                            "remixSuggestions": curr_remixes
                        }
                        
                        zip_buffer = create_final_zip(export_data, st.session_state.images)
                        st.download_button("⬇️ Download ZIP", data=zip_buffer.getvalue(), file_name="dataset.zip", mime="application/zip", type="primary", use_container_width=True)

            # 2. Scrollable Area
            st.markdown('<div class="right-scroll-area">', unsafe_allow_html=True)

            # Main Prompt
            st.markdown("#### 📝 Main Prompt")
            default_text = item['original_prompt_text']
            if not default_text.strip().lower().startswith("create"):
                default_text = "Create an image of " + default_text
            main_prompt = st.text_area("main_hidden", value=default_text, height=80, key=f"m_{current_id}", label_visibility="collapsed")

            st.markdown("---")

            # Batch Paste
            with st.expander("📋 Paste Remix Text (Replace)", expanded=False):
                st.text_area("Paste here", height=100, key="batch_input_area", label_visibility="collapsed", placeholder="Title\nCreate...")
                session_key = f"remix_{current_id}"
                st.button("Parse & Replace", on_click=batch_parse_callback, args=(session_key, current_id))
                if st.session_state.get("_parse_success"):
                    st.success(f"Updated {st.session_state['_parsed_count']} items!")
                    st.session_state["_parse_success"] = False
                if st.session_state.get("_parse_error"):
                    st.warning("No valid prompts found.")
                    st.session_state["_parse_error"] = False

            # Remix Cards
            st.markdown("#### 🎨 Remix Suggestions")
            if session_key not in st.session_state:
                st.session_state[session_key] = [get_random_remix() for _ in range(3)]
            current_remixes = st.session_state[session_key]

            # Horizontal Cards (3 Columns)
            r_cols = st.columns(3)
            for i, col in enumerate(r_cols):
                with col:
                    with st.container(border=True):
                        # Title
                        l_key = f"l_{current_id}_{i}"
                        if l_key not in st.session_state: st.session_state[l_key] = current_remixes[i]['label']
                        l_val = st.text_input(f"L{i}", value=current_remixes[i]['label'], key=l_key, label_visibility="collapsed", placeholder="Label")
                        
                        # Dice (Full width in card)
                        st.button("🎲", key=f"rnd_{current_id}_{i}", on_click=randomize_callback, args=(i, session_key, current_id), use_container_width=True)

                        # Prompt
                        p_key = f"p_{current_id}_{i}"
                        if p_key not in st.session_state: st.session_state[p_key] = current_remixes[i]['prompt']
                        p_val = st.text_area(f"P{i}", value=current_remixes[i]['prompt'], height=100, key=p_key, label_visibility="collapsed", placeholder="Prompt")

                        # Verify
                        if st.button("Verify", key=f"v_{current_id}_{i}", use_container_width=True):
                            clean = urllib.parse.quote(p_val)
                            seed = random.randint(0, 9999)
                            url = f"https://image.pollinations.ai/prompt/{clean}?seed={seed}&width=400&height=400&nologo=true"
                            st.session_state[f"poll_img_{current_id}_{i}"] = url
                        
                        if f"poll_img_{current_id}_{i}" in st.session_state:
                            st.image(st.session_state[f"poll_img_{current_id}_{i}"], use_container_width=True)

            st.markdown('</div>', unsafe_allow_html=True) # End scrollable

            # Bottom Bar (Outside scroll)
            st.markdown("<br>", unsafe_allow_html=True)
            b_col1, b_col2 = st.columns([1, 4])
            with b_col1:
                if st.button("⬅️", key="prev_bottom", use_container_width=True, disabled=st.session_state.current_idx == 0):
                    st.session_state.current_idx -= 1
                    st.rerun()
            with b_col2:
                if st.button("💾 Save & Next", type="primary", use_container_width=True):
                    final_json = { 
                        "id": current_id, 
                        "prompt": main_prompt, 
                        "remixSuggestions": current_remixes
                    }
                    st.session_state.processed_results[f"{current_id}.json"] = final_json
                    if st.session_state.current_idx < len(st.session_state.data) - 1:
                        st.session_state.current_idx += 1
                        st.rerun()
                    else:
                        st.balloons()
                        st.success("All Done! Check Export.")

# ================= TAB 2: JSON FIXER =================
with tab_fix:
    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 2, 1])
    with c2:
        st.markdown(f"<div style='text-align: center; margin-bottom:20px;'><h1>🔢 JSON ID Fixer</h1><p style='color:#64748B'>Upload dataset.json • Renumber IDs</p></div>", unsafe_allow_html=True)
        with st.container(border=True):
            up_json = st.file_uploader("Upload dataset.json", type=["json"])
            new_start = st.number_input("New Start ID", value=1001, step=1)
            
            if up_json:
                if st.button("🚀 Process & Renumber", type="primary", use_container_width=True):
                    new_json_str, error = renumber_json_ids(up_json, new_start)
                    if error:
                        st.error(error)
                    else:
                        st.success("IDs updated successfully!")
                        st.download_button("⬇️ Download New JSON", data=new_json_str, file_name="dataset_renumbered.json", mime="application/json", type="primary", use_container_width=True)

# ================= TAB 3: IMAGE EXTRACTOR =================
with tab_extract:
    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 2, 1])
    with c2:
        st.markdown(f"<div style='text-align: center; margin-bottom:20px;'><h1>🖼️ PPT Image Extractor</h1><p style='color:#64748B'>Extract all images & Rename by ID</p></div>", unsafe_allow_html=True)
        with st.container(border=True):
            ext_ppt = st.file_uploader("Upload Presentation (.pptx)", type=["pptx"], key="ext_uploader")
            ext_start = st.number_input("Start Filename ID", value=453, step=1, key="ext_start")
            
            if ext_ppt:
                if st.button("🚀 Extract & Zip", type="primary", use_container_width=True):
                    with st.spinner("Extracting..."):
                        zip_buf, count = extract_images_from_ppt(ext_ppt, ext_start)
                        if zip_buf:
                            st.success(f"Extracted {count} images!")
                            st.download_button("⬇️ Download Images ZIP", data=zip_buf.getvalue(), file_name="images_extracted.zip", mime="application/zip", type="primary", use_container_width=True)
                        else:
                            st.error(f"Error: {count}")
